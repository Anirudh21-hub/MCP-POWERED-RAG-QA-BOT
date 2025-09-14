# Standard library imports
import os
import requests
import re
import chardet
import logging
from typing import Optional, List, Dict, Any

# PDF processing library
import pdfplumber

# Data manipulation for CSV files
import pandas as pd
import numpy as np

# Markdown processing
import markdown
from markdown.extensions import codehilite, tables, toc

# PowerPoint file processing
from pptx import Presentation

# Word document processing
from docx import Document as DocxDocument

# Environment variable management
from dotenv import load_dotenv

# MCP server framework
from fastmcp import FastMCP

# LangChain core components
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_community.embeddings.spacy_embeddings import SpacyEmbeddings
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate
from langchain_community.vectorstores import FAISS as LangchainFAISS

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# Load environment variables from .env file
load_dotenv()
# Initialize MCP server instance
mcp = FastMCP("MCP Powered RAG Server")

# Global variables for vector store and document management
vectorstore = None  # FAISS vector store for document embeddings
doc_objects = []  # List to store all processed documents

# Initialize text embedding model using spaCy
embedder = SpacyEmbeddings(model_name="en_core_web_sm")
# Configure text splitter for chunking documents
text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)


def _clean_text(text: str) -> str:
    """Clean and normalize extracted text."""
    if not text:
        return ""

    # Remove excessive whitespace and normalize line breaks
    text = re.sub(r"\s+", " ", text)

    # Remove special characters that might interfere with processing
    text = re.sub(r"[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}\"\'\/\\]", "", text)

    # Normalize quotes
    text = re.sub(r"[\u201c\u201d]", '"', text)
    text = re.sub(r"[\u2018\u2019]", "'", text)

    # Remove excessive punctuation
    text = re.sub(r"\.{3,}", "...", text)
    text = re.sub(r"\!{2,}", "!", text)
    text = re.sub(r"\?{2,}", "?", text)

    # Clean up common PDF artifacts
    text = re.sub(r"\f", "\n", text)  # Form feed to newline
    text = re.sub(r"\x00", "", text)  # Remove null characters
    text = re.sub(
        r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text
    )  # Remove control characters

    # Remove page numbers and headers/footers (common patterns)
    text = re.sub(r"^\s*\d+\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^Page \d+ of \d+$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*-\s*\d+\s*-\s*$", "", text, flags=re.MULTILINE)

    # Clean up multiple consecutive newlines
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def _detect_encoding(file_path: str) -> str:
    """Detect file encoding using chardet."""
    try:
        with open(file_path, "rb") as f:
            raw_data = f.read(10000)  # Read first 10KB for detection
            result = chardet.detect(raw_data)
            encoding = result.get("encoding", "utf-8")
            confidence = result.get("confidence", 0)

            # Fallback to utf-8 if confidence is low
            if confidence < 0.7:
                encoding = "utf-8"

            logger.info(f"Detected encoding: {encoding} (confidence: {confidence:.2f})")
            return encoding
    except Exception as e:
        logger.warning(f"Failed to detect encoding for {file_path}: {e}")
        return "utf-8"


def _extract_text_from_pdf(file_path):
    """Extract text from PDF files with improved parsing."""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            logger.info(f"Processing PDF with {len(pdf.pages)} pages")

            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    # Try multiple extraction methods for better results
                    page_text = page.extract_text()

                    # If no text found, try extracting from tables
                    if not page_text or len(page_text.strip()) < 10:
                        tables = page.extract_tables()
                        if tables:
                            table_text = ""
                            for table in tables:
                                for row in table:
                                    if row:
                                        table_text += (
                                            " ".join([cell for cell in row if cell])
                                            + "\n"
                                        )
                            if table_text.strip():
                                page_text = table_text

                    # If still no text, try extracting from text objects
                    if not page_text or len(page_text.strip()) < 10:
                        chars = page.chars
                        if chars:
                            page_text = "".join([char["text"] for char in chars])

                    if page_text:
                        # Clean the page text
                        cleaned_text = _clean_text(page_text)
                        if cleaned_text:
                            text += cleaned_text + "\n\n"
                            logger.debug(
                                f"Extracted {len(cleaned_text)} characters from page {page_num}"
                            )
                        else:
                            logger.warning(f"No clean text found on page {page_num}")
                    else:
                        logger.warning(f"No text found on page {page_num}")

                except Exception as e:
                    logger.error(f"Error processing page {page_num}: {e}")
                    continue

    except Exception as e:
        logger.error(f"Error opening PDF file {file_path}: {e}")
        raise

    if not text.strip():
        logger.warning(f"No text extracted from PDF: {file_path}")

    return _clean_text(text)


def _extract_text_from_pptx(file_path):
    """Extract text from PowerPoint files with improved parsing."""
    text = ""
    try:
        prs = Presentation(file_path)
        logger.info(f"Processing PowerPoint with {len(prs.slides)} slides")

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""

            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                    elif hasattr(shape, "text_frame"):
                        # Handle text frames (more complex text layouts)
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_text += paragraph.text.strip() + "\n"
                except Exception as e:
                    logger.debug(f"Error processing shape in slide {slide_num}: {e}")
                    continue

            if slide_text.strip():
                text += f"--- Slide {slide_num} ---\n{slide_text}\n"
                logger.debug(f"Extracted text from slide {slide_num}")
            else:
                logger.warning(f"No text found in slide {slide_num}")

        if not text.strip():
            logger.warning(f"No text extracted from PowerPoint: {file_path}")

    except Exception as e:
        logger.error(f"Error processing PowerPoint file {file_path}: {e}")
        raise

    return _clean_text(text)


def _extract_text_from_docx(file_path):
    """Extract text from Word documents with improved parsing."""
    text = ""
    try:
        doc = DocxDocument(file_path)
        logger.info(f"Processing Word document with {len(doc.paragraphs)} paragraphs")

        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text.strip() + "\n"

        # Extract text from tables
        for table in doc.tables:
            table_text = ""
            for row in table.rows:
                row_text = ""
                for cell in row.cells:
                    if cell.text.strip():
                        row_text += cell.text.strip() + " | "
                if row_text:
                    table_text += row_text.rstrip(" | ") + "\n"
            if table_text.strip():
                text += f"\n--- Table ---\n{table_text}\n"

        # Extract text from headers and footers
        for section in doc.sections:
            if section.header:
                header_text = ""
                for para in section.header.paragraphs:
                    if para.text.strip():
                        header_text += para.text.strip() + "\n"
                if header_text.strip():
                    text += f"\n--- Header ---\n{header_text}\n"

            if section.footer:
                footer_text = ""
                for para in section.footer.paragraphs:
                    if para.text.strip():
                        footer_text += para.text.strip() + "\n"
                if footer_text.strip():
                    text += f"\n--- Footer ---\n{footer_text}\n"

        if not text.strip():
            logger.warning(f"No text extracted from Word document: {file_path}")

    except Exception as e:
        logger.error(f"Error processing Word document {file_path}: {e}")
        raise

    return _clean_text(text)


def _extract_text_from_csv(file_path):
    """Extract text from CSV files with improved parsing."""
    text = ""
    try:
        # Detect encoding first
        encoding = _detect_encoding(file_path)

        # Try different CSV parsing options
        for sep in [",", ";", "\t", "|"]:
            try:
                df = pd.read_csv(
                    file_path, encoding=encoding, sep=sep, on_bad_lines="skip"
                )
                if len(df.columns) > 1:  # Valid CSV with multiple columns
                    logger.info(
                        f"Successfully parsed CSV with separator '{sep}' and {len(df)} rows"
                    )
                    break
            except Exception:
                continue
        else:
            # Fallback to default comma separator
            df = pd.read_csv(file_path, encoding=encoding, on_bad_lines="skip")

        # Convert DataFrame to structured text
        if not df.empty:
            # Add column headers
            text += "Columns: " + ", ".join(df.columns.astype(str)) + "\n\n"

            # Add data rows with better formatting
            for idx, row in df.iterrows():
                row_text = ""
                for col in df.columns:
                    value = str(row[col]) if pd.notna(row[col]) else ""
                    row_text += f"{col}: {value}; "
                text += row_text.rstrip("; ") + "\n"

            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text += "\nSummary Statistics:\n"
                for col in numeric_cols:
                    stats = df[col].describe()
                    text += f"{col}: mean={stats['mean']:.2f}, std={stats['std']:.2f}, min={stats['min']:.2f}, max={stats['max']:.2f}\n"

        if not text.strip():
            logger.warning(f"No data found in CSV file: {file_path}")

    except Exception as e:
        logger.error(f"Error reading CSV file {file_path}: {e}")
        # Fallback: read as plain text
        try:
            encoding = _detect_encoding(file_path)
            with open(file_path, "r", encoding=encoding) as f:
                text = f.read()
            logger.info(f"Fallback: read CSV as plain text")
        except Exception as fallback_error:
            logger.error(f"Fallback also failed: {fallback_error}")
            text = f"Error: Could not parse CSV file. Original error: {e}"

    return _clean_text(text)


def _extract_text_from_txt(file_path):
    """Extract text from plain text and markdown files with improved parsing."""
    text = ""
    try:
        # Detect encoding first
        encoding = _detect_encoding(file_path)

        with open(file_path, "r", encoding=encoding) as f:
            content = f.read()

        # Check if it's a markdown file
        if file_path.lower().endswith((".md", ".markdown")):
            # Use markdown extensions for better parsing
            md = markdown.Markdown(
                extensions=["codehilite", "tables", "toc", "fenced_code", "nl2br"]
            )
            html = md.convert(content)

            # Better HTML to text conversion
            # Remove script and style elements
            html = re.sub(
                r"<(script|style)[^>]*>.*?</\1>",
                "",
                html,
                flags=re.DOTALL | re.IGNORECASE,
            )

            # Convert HTML entities
            html = html.replace("&lt;", "<").replace("&gt;", ">").replace("&amp;", "&")
            html = (
                html.replace("&quot;", '"').replace("&#39;", "'").replace("&nbsp;", " ")
            )

            # Remove HTML tags but preserve structure
            text = re.sub(r"<h[1-6][^>]*>", "\n# ", html)  # Headers
            text = re.sub(r"<p[^>]*>", "\n", text)  # Paragraphs
            text = re.sub(r"<br[^>]*>", "\n", text)  # Line breaks
            text = re.sub(r"<li[^>]*>", "\n• ", text)  # List items
            text = re.sub(r"<blockquote[^>]*>", "\n> ", text)  # Blockquotes
            text = re.sub(r"<code[^>]*>", "`", text)  # Inline code
            text = re.sub(r"<pre[^>]*>", "\n```\n", text)  # Code blocks
            text = re.sub(r"</pre>", "\n```\n", text)
            text = re.sub(r"<[^>]+>", "", text)  # Remove remaining HTML tags

            # Clean up the text
            text = re.sub(r"\n{3,}", "\n\n", text)  # Remove excessive newlines
            text = re.sub(
                r"^\s+", "", text, flags=re.MULTILINE
            )  # Remove leading whitespace

        else:
            text = content  # Use content as-is for plain text files

        logger.info(f"Extracted {len(text)} characters from text file: {file_path}")

    except Exception as e:
        logger.error(f"Error reading text file {file_path}: {e}")
        # Fallback to basic reading
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as fallback_error:
            logger.error(f"Fallback also failed: {fallback_error}")
            text = f"Error: Could not read text file. Original error: {e}"

    return _clean_text(text)


def _extract_text_from_file(file_path):
    """Extract text from various file formats based on file extension with improved error handling."""
    file_ext = os.path.splitext(file_path)[1].lower()

    logger.info(f"Processing file: {file_path} (extension: {file_ext})")

    try:
        if file_ext == ".pdf":
            return _extract_text_from_pdf(file_path)
        elif file_ext == ".pptx":
            return _extract_text_from_pptx(file_path)
        elif file_ext == ".docx":
            return _extract_text_from_docx(file_path)
        elif file_ext == ".csv":
            return _extract_text_from_csv(file_path)
        elif file_ext in [".txt", ".md", ".markdown"]:
            return _extract_text_from_txt(file_path)
        else:
            raise ValueError(
                f"Unsupported file format: {file_ext}. Supported formats: .pdf, .pptx, .docx, .csv, .txt, .md, .markdown"
            )
    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
        raise


def _get_supported_extensions():
    """Get list of supported file extensions."""
    return [
        ".pdf",
        ".pptx",
        ".docx",
        ".csv",
        ".txt",
        ".md",
        ".markdown",
    ]  # Return list of supported file extensions


def get_embeddings(file_path):
    """Extract text from file and create embeddings with improved error handling."""
    global vectorstore, doc_objects

    try:
        logger.info(f"Starting to process file: {file_path}")

        # Extract text from file
        text = _extract_text_from_file(file_path)

        if not text or not text.strip():
            logger.warning(f"No text extracted from {file_path}")
            return False

        logger.info(f"Extracted {len(text)} characters from {file_path}")

        # Split text into chunks
        chunks = text_splitter.split_text(text)
        logger.info(f"Split text into {len(chunks)} chunks")

        # Create Document objects with metadata
        docs = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():  # Only add non-empty chunks
                doc = Document(
                    page_content=chunk,
                    metadata={
                        "source": file_path,
                        "chunk_index": i,
                        "file_type": os.path.splitext(file_path)[1].lower(),
                        "total_chunks": len(chunks),
                    },
                )
                docs.append(doc)

        if not docs:
            logger.warning(f"No valid chunks created from {file_path}")
            return False

        # Add documents to global list
        doc_objects.extend(docs)

        # Create or update vector store
        if vectorstore is None:
            logger.info("Creating new vector store")
            vectorstore = LangchainFAISS.from_documents(docs, embedder)
        else:
            logger.info("Adding documents to existing vector store")
            vectorstore.add_documents(docs)

        file_ext = os.path.splitext(file_path)[1].upper()
        logger.info(
            f"Successfully indexed {len(docs)} chunks from {file_ext} file: {os.path.basename(file_path)}"
        )

        return True

    except Exception as e:
        logger.error(f"Failed to process {file_path}: {str(e)}", exc_info=True)
        raise


# Custom prompt template for RAG responses
custom_prompt = PromptTemplate.from_template(
    """You are a AI Assistant provides answer to the user question and follow the instructions. 
Context: {context} 
Question: {question} 
History: {chat_history} 
."""
)


def get_llm():
    api_key = os.environ.get("GROQ_API_KEY")  # Get API key from environment
    if not api_key:  # Check if API key is not set
        raise RuntimeError(  # Raise error if API key is missing
            "GROQ_API_KEY is not set. Set it in your environment or .env file."
        )
    model = os.environ.get(
        "GROQ_MODEL", "llama-3.3-70b-versatile"
    )  # Get model name with default
    return ChatGroq(api_key=api_key, model=model)  # Return configured ChatGroq instance


def _create_chain():
    if vectorstore is None:  # Check if vector store is not initialized
        raise ValueError(
            "No documents indexed. Please index a file first."
        )  # Raise error if no documents

    retriever = vectorstore.as_retriever(
        search_kwargs={"k": 5}
    )  # Create retriever with top 5 results

    return (
        ConversationalRetrievalChain.from_llm(  # Create conversational retrieval chain
            llm=get_llm(),  # Use configured LLM
            retriever=retriever,  # Use document retriever
            memory=ConversationBufferMemory(  # Configure conversation memory
                memory_key="chat_history", output_key="answer", return_messages=True
            ),
            combine_docs_chain_kwargs={
                "prompt": custom_prompt
            },  # Use custom prompt template
            output_key="answer",  # Set output key for answers
        )
    )


@mcp.tool()  # Register as MCP tool
def rag_query(question: str) -> str:
    print("[debug] RAG Query:", question)  # Print debug information
    chain = _create_chain()  # Create retrieval chain
    result = chain.invoke({"question": question})  # Invoke chain with question
    response = result["answer"]  # Extract answer from result

    if not response:  # Check if response is empty
        result = chain.invoke({"question": question})  # Retry with same question
        response = result["answer"]  # Extract answer from retry result

    return response  # Return final response


@mcp.tool()  # Register as MCP tool
def index_file(file_path: str) -> str:
    """Index a file for RAG. Supports PDF, PPTX, DOCX, CSV, TXT, and Markdown files."""
    logger.info(f"Indexing file: {file_path}")

    # Check if file exists
    if not os.path.exists(file_path):
        error_msg = f"Error: File not found: {file_path}"
        logger.error(error_msg)
        return error_msg

    # Check file extension
    file_ext = os.path.splitext(file_path)[1].lower()
    supported_exts = _get_supported_extensions()

    if file_ext not in supported_exts:
        error_msg = f"Error: Unsupported file format. Supported formats: {', '.join(supported_exts)}"
        logger.error(error_msg)
        return error_msg

    try:
        success = get_embeddings(file_path)
        if success:
            success_msg = f"Successfully indexed {os.path.basename(file_path)}"
            logger.info(success_msg)
            return success_msg
        else:
            error_msg = (
                f"Failed to index {os.path.basename(file_path)}: No text extracted"
            )
            logger.warning(error_msg)
            return error_msg
    except Exception as e:
        error_msg = f"Error indexing file: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return error_msg


@mcp.tool()  # Register as MCP tool
def get_supported_formats() -> str:
    """Get list of supported file formats."""
    formats = _get_supported_extensions()  # Get supported extensions
    return f"Supported file formats: {', '.join(formats)}"  # Return formatted string


@mcp.tool()  # Register as MCP tool
def index_pdf_file(file_path: str) -> str:
    """Legacy function for backward compatibility. Use index_file instead."""
    return index_file(file_path)  # Call the new index_file function


if __name__ == "__main__":  # Check if script is run directly
    sample_dir = os.path.join(
        os.path.dirname(__file__), "sample_files"
    )  # Get sample files directory
    if os.path.exists(sample_dir):  # Check if sample directory exists
        supported_exts = _get_supported_extensions()  # Get supported file extensions
        for file in os.listdir(sample_dir):  # Iterate through files in sample directory
            file_path = os.path.join(sample_dir, file)  # Create full file path
            file_ext = os.path.splitext(file)[1].lower()  # Get file extension
            if file_ext in supported_exts:  # Check if file type is supported
                try:
                    get_embeddings(file_path)  # Process file and create embeddings
                except Exception as e:
                    print(
                        f"[warning] Failed to index {file}: {e}"
                    )  # Print warning for failed files

    port = int(
        os.environ.get("PORT", "8000")
    )  # Get port from environment or use default
    host = os.environ.get(
        "HOST", "127.0.0.1"
    )  # Get host from environment or use default
    mcp.run(
        transport="sse", host=host, port=port
    )  # Start MCP server with SSE transport
